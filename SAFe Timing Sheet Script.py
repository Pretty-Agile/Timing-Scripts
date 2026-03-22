#!/usr/bin/env python3
"""
Pretty Agile — SAFe Timing v4.6
- Interactive configuration on each run (skip with --no-config):
    · Day timing: 9:00-18:00, 8:30-18:00, or custom start time
      (lunch set 4 h after start; breaks auto-derived hourly)
    · Lesson balancing: weighted (natural overflow), even (modules split
      across N days), or custom (MODULE:DAY pin pairs, e.g. 7:4)
- Hard daily cut-off: 9:00→18:00 (default); 8:30→18:00 (optional)
- Visible 'Start of Day' row (0 mins, grey) resets time each day
- Do NOT skip breaks/lunch: if a Break/Lunch would overflow today, end day and place it first thing next day
- Final day: if time remains, add 'Close – Photo, Feedback & Exam' to fill remaining time (formula)
- 12-hour h:mm AM/PM formatting; grey shading for Break/Lunch/Start of Day
- Excel formulas for Minutes/Total/Start/End; no Day/Deck columns
- Fixed breaks (10m): 8:30→ 09:30,10:30,11:30,14:30,15:30,16:30,17:30; 9:00→ 10:00,11:00,12:00,15:00,16:00,17:00
- Breaks and lunch fire within ±5 min of target; lessons split mid-lesson if target falls inside them
- Split format: "4.2 Lesson (up to 4.14)" / "4.2 Lesson (from 4.14)" — activity minutes placed in whichever half contains their slide

"""
import argparse, re, os, sys, platform, time
from pathlib import Path
from typing import List, Tuple, Optional, Dict, Any
from concurrent.futures import ProcessPoolExecutor, as_completed
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from datetime import timedelta
from openpyxl import Workbook
from openpyxl.styles import PatternFill

# ---------- regex ----------
ACTIVITY_PREFIX = re.compile(r'^\s*(activity|discussion|video|action\s*plan)\b', re.I)
VIDEO_PREFIX    = re.compile(r'^\s*video\b', re.I)
PI_PLANNING_DECK = re.compile(r'leading\s+safe|safe\s+scrum\s+master|implementing\s+safe', re.I)
PI_PLANNING_SLIDE = re.compile(r'pi\s+planning', re.I)
SUBHEADER_RE = re.compile(r'^\s*\d+\.\d+\b(?!\.)')
DUR_PATTERNS = [
    (re.compile(r'(\d{1,3})\s*-\s*(\d{1,3})\s*m(?:in(?:s|utes)?)?\b', re.I),
     lambda m: max(int(m.group(1)), int(m.group(2)))),
    (re.compile(r'(\d{1,2})\s*h(?:ours?)?\s*(\d{1,3})\s*m(?:in(?:s|utes)?)?\b', re.I),
     lambda m: int(m.group(1)) * 60 + int(m.group(2))),
    (re.compile(r'(\d{1,2})\s*h(?:ours?)?\b', re.I),
     lambda m: int(m.group(1)) * 60),
    (re.compile(r'(\d{1,3})\s*m(?:in(?:s|utes)?)?\b', re.I),
     lambda m: int(m.group(1))),
    (re.compile(r'(?:^|\D)(\d{1,3})(?!\s*[\.\d])(?:\D|$)'),
     lambda m: int(m.group(1))),
]

# ---------- utils ----------
def clean_one_line(s: str) -> str:
    return " ".join((s or "").split())

def strip_leading_two_digits(name: str) -> str:
    return re.sub(r'^\s*\d{2}\s+', '', name).strip()

def deck_label_from_filename(fname: str) -> str:
    base = Path(fname).stem
    base = re.sub(r'\s*\(Converted\)\s*$', '', base, flags=re.I)
    base = re.sub(r'\s*\(\d+(?:\.\d+)*\)\s*$', '', base)
    return base.strip()

def slide_title(slide) -> str:
    for shp in slide.shapes:
        try:
            if shp.is_placeholder and shp.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                if shp.has_text_frame:
                    t = clean_one_line(shp.text or "")
                    if t:
                        return t
        except Exception:
            pass
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            t = clean_one_line(shp.text or "")
            if t:
                return t
    return ""

def slide_all_text(slide) -> str:
    parts = []
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            if shp.text:
                parts.append(shp.text)
    try:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            parts.append(slide.notes_slide.notes_text_frame.text or "")
    except Exception:
        pass
    return "\n".join(p.strip() for p in parts if p and p.strip())

def is_subheader_title(title: str) -> bool:
    return bool(SUBHEADER_RE.match(title or ""))

def normalized_sublesson_name(slide, header_title: str, precomputed_body: str = None) -> str:
    title = header_title or ""
    code_m = re.match(r'^\s*(\d+\.\d+)\b', title)
    code = code_m.group(1) if code_m else ""
    m = re.match(r'^\s*\d+\.\d+\s*[:\-–]?\s*(.+)$', title)
    if m and m.group(1).strip():
        return f"{code} {m.group(1).strip()}"
    body = precomputed_body if precomputed_body is not None else slide_all_text(slide)
    for line in (l.strip() for l in body.splitlines() if l.strip()):
        m2 = re.match(r'^\s*\d+\.\d+\s*[:\-–]?\s*(.+)$', line)
        if m2 and m2.group(1).strip():
            return f"{code} {m2.group(1).strip()}"
    lines = [l.strip() for l in body.splitlines() if l.strip()]
    longest = max(lines, key=len) if lines else "Sub-lesson"
    return f"{code} {longest}"

def is_activity_slide(title: str) -> bool:
    return bool(ACTIVITY_PREFIX.match(title or ""))

def is_video_slide(title: str) -> bool:
    return bool(VIDEO_PREFIX.match(title or ""))

def extract_minutes(text: str) -> Optional[int]:
    s = (text or "").replace("\u2013", "-").replace("\u2014", "-")
    for pat, func in DUR_PATTERNS:
        m = pat.search(s)
        if m:
            try:
                return max(0, func(m))
            except Exception:
                continue
    return None

# ---------- parse deck ----------
def parse_deck(path: Path, deck_label: str) -> List[Tuple[str, int, int, str]]:
    prs = Presentation(str(path))
    slides = list(prs.slides)
    titles = [slide_title(s) for s in slides]
    texts = [slide_all_text(s) for s in slides]
    combined = [(titles[i] + "\n" + texts[i]).strip() for i in range(len(slides))]

    header_idxs, header_names = [], {}
    for i, t in enumerate(titles):
        if is_subheader_title(t):
            header_idxs.append(i)
            header_names[i] = normalized_sublesson_name(slides[i], t, texts[i])

    deck_clean_no_lead = strip_leading_two_digits(deck_label)
    intro_label = "1.0 " + deck_clean_no_lead
    if header_idxs:
        first_name = header_names[header_idxs[0]] or titles[header_idxs[0]]
        m = re.match(r'^\s*(\d+)\.\d+\b', first_name or "")
        if m:
            intro_label = f"{m.group(1)}.0 {deck_clean_no_lead}"
    else:
        m = re.match(r'^\s*(\d+)', deck_label)
        if m:
            intro_label = f"{int(m.group(1))}.0 {deck_clean_no_lead}"

    is_pi_deck = bool(PI_PLANNING_DECK.search(deck_label))

    def sum_activity(start: int, end: int) -> int:
        total = 0
        for i in range(start, end):
            if is_activity_slide(titles[i]):
                mins = extract_minutes(combined[i])
                if mins is not None:
                    if is_pi_deck and mins == 50 and PI_PLANNING_SLIDE.search(combined[i]):
                        mins = 30
                    total += mins
        return total

    def count_video(start: int, end: int) -> int:
        return sum(1 for i in range(start, end) if is_video_slide(titles[i]))

    def adjusted_row(name, start, end):
        """Return (name, slide_count, activity_mins, kind, slide_start, acts_per_slide) with video slides contributing 2 mins each."""
        vid = count_video(start, end)
        count = max(0, end - start) - vid
        act = sum_activity(start, end) + vid * 2  # 2 mins per video slide replaces mins_per_slide
        # Build per-counted-slide activity list; video mins accumulate to the next counted slide
        acts: List[int] = []
        pending_vid = 0
        for i in range(start, end):
            if is_video_slide(titles[i]):
                pending_vid += 2
            else:
                slide_act = pending_vid
                pending_vid = 0
                if is_activity_slide(titles[i]):
                    mins = extract_minutes(combined[i])
                    if mins is not None:
                        if is_pi_deck and mins == 50 and PI_PLANNING_SLIDE.search(combined[i]):
                            mins = 30
                        slide_act += mins
                acts.append(slide_act)
        if pending_vid and acts:
            acts[-1] += pending_vid  # trailing video slides → last counted slide
        return [name, count, act, "lesson", start, acts]

    rows: List[Tuple[str, int, int, str]] = []
    if not header_idxs:
        rows.append(adjusted_row(intro_label, 0, len(slides)))
        return rows

    first = header_idxs[0]
    if first > 0:
        rows.append(adjusted_row(intro_label, 0, first))

    for k, start in enumerate(header_idxs):
        end = header_idxs[k + 1] if k + 1 < len(header_idxs) else len(slides)
        name = header_names.get(start) or f"Sub-lesson {k+1}"
        rows.append(adjusted_row(name, start, end))

    return rows

# ---------- time helpers ----------
def parse_hhmm(s: str) -> timedelta:
    h, m = s.split(":")
    return timedelta(hours=int(h), minutes=int(m))

def fixed_break_targets(day_window: str) -> List[timedelta]:
    if day_window == "8:30-18:00":
        return [
            timedelta(hours=9, minutes=30),
            timedelta(hours=10, minutes=30),
            timedelta(hours=11, minutes=30),
            timedelta(hours=14, minutes=30),
            timedelta(hours=15, minutes=30),
            timedelta(hours=16, minutes=30),
            timedelta(hours=17, minutes=30),
        ]
    else:
        return [
            timedelta(hours=10, minutes=0),
            timedelta(hours=11, minutes=0),
            timedelta(hours=12, minutes=0),
            timedelta(hours=15, minutes=0),
            timedelta(hours=16, minutes=0),
            timedelta(hours=17, minutes=0),
        ]

def compute_lunch_time(start_str: str) -> str:
    """Return a lunch time string 4 hours after start (HH:MM)."""
    lunch = parse_hhmm(start_str) + timedelta(hours=4)
    total_mins = int(lunch.total_seconds() // 60)
    return f"{total_mins // 60:02d}:{total_mins % 60:02d}"


def compute_break_targets(start_str: str, end_str: str, lunch_str: str) -> List[timedelta]:
    """Compute hourly break targets for a custom day window.

    Generates one break per hour starting at start+1h up to end, excluding
    the lunch slot and the hour immediately after lunch.
    """
    start = parse_hhmm(start_str)
    end = parse_hhmm(end_str)
    lunch = parse_hhmm(lunch_str)
    post_lunch = lunch + timedelta(hours=1)
    targets: List[timedelta] = []
    t = start + timedelta(hours=1)
    while t < end:
        if t != lunch and t != post_lunch:
            targets.append(t)
        t += timedelta(hours=1)
    return targets


# ---------- sequence with hard cut-off ----------
def build_sequence(
    rows_by_deck: List[Dict[str, Any]],
    day_start: str,
    day_end: str,
    lunch_time: str,
    day_window: str,
    mins_per_slide: float,
    lesson_day_map: Dict[str, int] = None,
    break_targets: List[timedelta] = None,
) -> List[Dict[str, Any]]:
    """
    Returns ordered rows (lessons + Start of Day + breaks + lunch + final close),
    respecting daily cut-offs and 'do not skip' for breaks/lunch.

    Breaks and lunch fire within ±5 minutes of their target times.  If a target
    falls well inside a lesson (neither lesson boundary is within the 5-minute
    window), the lesson is split at the target time:
      "4.2 Lesson Name (up to 4.14)"  ← first half  (all activity minutes here)
      "4.2 Lesson Name (from 4.14)"   ← second half (0 activity minutes)
    where 4.14 is the module.slide notation for the first slide of the second half.
    """
    BREAK_WINDOW = timedelta(minutes=5)

    start_td = parse_hhmm(day_start)
    end_td = parse_hhmm(day_end)
    lunch_td = parse_hhmm(lunch_time)
    lunch_earliest = lunch_td - timedelta(minutes=15)
    targets = break_targets if break_targets is not None else fixed_break_targets(day_window)

    # Make a mutable copy so we can replace rows in-place during splits
    rows_by_deck = list(rows_by_deck)

    out: List[Dict[str, Any]] = []
    day = 1
    t = start_td
    lunch_done_today = False
    next_break_idx = 0
    first_lesson_scheduled_today = False

    def start_day_row():
        nonlocal t, lunch_done_today, next_break_idx, first_lesson_scheduled_today
        t = start_td
        lunch_done_today = False
        next_break_idx = 0
        first_lesson_scheduled_today = False
        out.append({"Sub-lesson": "Start of Day", "Slides": 0, "Activity Minutes": 0, "kind": "sod"})

    def end_day_and_start_next():
        nonlocal day
        day += 1
        start_day_row()

    # Begin Day 1
    start_day_row()

    # Helper to append Break/Lunch respecting overflow (move to next day if needed)
    def append_block(kind: str, minutes: int):
        nonlocal t
        if kind == "break" and out and out[-1]["kind"] == "break":
            return
        if t + timedelta(minutes=minutes) > end_td:
            end_day_and_start_next()
        out.append({
            "Sub-lesson": "Break" if kind == "break" else "Lunch",
            "Slides": 0,
            "Activity Minutes": minutes,
            "kind": kind,
        })
        t += timedelta(minutes=minutes)

    def skip_post_lunch_break():
        nonlocal next_break_idx
        skip_target = (
            timedelta(hours=13, minutes=30) if day_window == "8:30-18:00"
            else timedelta(hours=14, minutes=0)
        )
        if next_break_idx < len(targets) and targets[next_break_idx] == skip_target:
            next_break_idx += 1

    def make_slide_ref(row: Dict, slides_before: int) -> str:
        """Return 'module.NN' for the first slide of the second half (1-indexed in deck)."""
        slide_num = row.get("slide_start", 0) + slides_before + 1
        deck_module = row.get("deck_module", "")
        return f"{deck_module}.{slide_num:02d}" if deck_module else str(slide_num)

    def compute_split(t_now: timedelta, target: timedelta,
                      slides: int, acts_per_slide: List[int]) -> Optional[Tuple[int, int]]:
        """
        Return (slides_1, slides_2) for splitting at target, or None if the split
        would leave one half empty (caller should treat as before/after instead).
        Iterates slide-by-slide so activity minutes on specific slides land in the
        correct half.
        """
        time_to_break = (target - t_now).total_seconds() / 60
        cumulative = 0.0
        slides_1 = 0
        for j in range(slides):
            slide_cost = mins_per_slide + (acts_per_slide[j] if j < len(acts_per_slide) else 0)
            if cumulative + slide_cost > time_to_break:
                break
            cumulative += slide_cost
            slides_1 = j + 1
        slides_2 = slides - slides_1
        if slides_1 == 0:
            return None   # nothing in first half → fire before instead
        if slides_2 == 0:
            return None   # nothing in second half → fire after instead
        return slides_1, slides_2

    _SPLIT_SUFFIX = re.compile(r'\s*\((?:from|up to)\s+[\d.]+\)\s*$')

    def _base_name(row: Dict) -> str:
        """Strip any existing split suffix so re-splits don't stack labels."""
        return _SPLIT_SUFFIX.sub('', row['Sub-lesson']).rstrip()

    def append_first_half(row: Dict, slides_1: int, activity_mins: int, ref: str):
        nonlocal t, first_lesson_scheduled_today
        half_mins = int(round(slides_1 * mins_per_slide)) + activity_mins
        out.append({
            "Sub-lesson": f"{_base_name(row)} (up to {ref})",
            "Slides": slides_1,
            "Activity Minutes": activity_mins,
            "kind": "lesson",
        })
        t += timedelta(minutes=half_mins)
        first_lesson_scheduled_today = True

    def make_second_half_row(row: Dict, slides_1: int, slides_2: int, ref: str,
                             acts_2: List[int]) -> Dict:
        return {
            "Deck": row.get("Deck", ""),
            "Sub-lesson": f"{_base_name(row)} (from {ref})",
            "Slides": slides_2,
            "Activity Minutes": sum(acts_2),
            "kind": "lesson",
            "slide_start": row.get("slide_start", 0) + slides_1,
            "deck_module": row.get("deck_module", ""),
            "acts_per_slide": acts_2,
        }

    # ---- Main scheduling loop ----
    i = 0
    current_module = ""
    while i < len(rows_by_deck):
        row = rows_by_deck[i]

        # ---- Force day advance when a new module starts (lesson_day_map) ----
        if lesson_day_map:
            module = row.get("deck_module", "")
            if module and module != current_module and module in lesson_day_map:
                target_day = lesson_day_map[module]
                while day < target_day:
                    end_day_and_start_next()
            current_module = module

        slides = int(row["Slides"])
        minutes = int(round(slides * mins_per_slide))
        activity_mins = int(row["Activity Minutes"])
        total = minutes + activity_mins
        acts_per_slide: List[int] = row.get("acts_per_slide") or [0] * slides

        # ---- Catch-up: fire any break targets already >5 min in the past ----
        while (next_break_idx < len(targets)
               and t > targets[next_break_idx] + BREAK_WINDOW):
            if first_lesson_scheduled_today:
                if t + timedelta(minutes=10 + total) <= end_td:
                    append_block("break", 10)
            next_break_idx += 1

        # ---- Lunch catch-up: past the window, fire before this lesson ----
        if not lunch_done_today and t >= lunch_earliest and t > lunch_td + BREAK_WINDOW:
            append_block("lunch", 45)
            lunch_done_today = True
            skip_post_lunch_break()

        # ---- Lunch: BEFORE or SPLIT ----
        if not lunch_done_today and t >= lunch_earliest:
            lesson_end = t + timedelta(minutes=total)
            if abs(t - lunch_td) <= BREAK_WINDOW:
                # Lesson start within ±5 min of lunch → LUNCH BEFORE
                append_block("lunch", 45)
                lunch_done_today = True
                skip_post_lunch_break()
            elif t < lunch_td and lesson_end > lunch_td + BREAK_WINDOW:
                # Lunch falls well inside the lesson → SPLIT at lunch_td
                split = compute_split(t, lunch_td, slides, acts_per_slide)
                if split is not None:
                    slides_1, slides_2 = split
                    ref = make_slide_ref(row, slides_1)
                    append_first_half(row, slides_1, sum(acts_per_slide[:slides_1]), ref)
                    append_block("lunch", 45)
                    lunch_done_today = True
                    skip_post_lunch_break()
                    rows_by_deck[i] = make_second_half_row(row, slides_1, slides_2, ref, acts_per_slide[slides_1:])
                    continue
                else:
                    # Degenerate: nothing fits before lunch → fire BEFORE
                    time_to_lunch = (lunch_td - t).total_seconds() / 60
                    if time_to_lunch <= activity_mins:
                        append_block("lunch", 45)
                        lunch_done_today = True
                        skip_post_lunch_break()
                    # else slides_2==0 → caught by LUNCH AFTER below

        # ---- Break: BEFORE or SPLIT ----
        if next_break_idx < len(targets):
            target = targets[next_break_idx]
            lesson_end = t + timedelta(minutes=total)
            if abs(t - target) <= BREAK_WINDOW:
                # Lesson start within ±5 min of target → BREAK BEFORE
                if first_lesson_scheduled_today:
                    if t + timedelta(minutes=10 + total) <= end_td:
                        append_block("break", 10)
                next_break_idx += 1
            elif t < target and lesson_end > target + BREAK_WINDOW:
                # Target falls well inside the lesson → SPLIT
                split = compute_split(t, target, slides, acts_per_slide)
                if split is not None:
                    slides_1, slides_2 = split
                    ref = make_slide_ref(row, slides_1)
                    append_first_half(row, slides_1, sum(acts_per_slide[:slides_1]), ref)
                    append_block("break", 10)
                    next_break_idx += 1
                    rows_by_deck[i] = make_second_half_row(row, slides_1, slides_2, ref, acts_per_slide[slides_1:])
                    continue
                else:
                    # Degenerate: check which side is closer
                    time_to_target = (target - t).total_seconds() / 60
                    if time_to_target <= activity_mins:
                        # All activity fills time before target → fire BEFORE
                        if first_lesson_scheduled_today:
                            if t + timedelta(minutes=10 + total) <= end_td:
                                append_block("break", 10)
                        next_break_idx += 1
                    # else slides_2==0 → caught by BREAK AFTER below

        # ---- Day overflow: split at end-of-day boundary or push whole lesson ----
        if t + timedelta(minutes=total) > end_td:
            remaining_mins = (end_td - t).total_seconds() / 60
            # Find how many slides (with their own activity) fit before end of day
            cumulative = 0.0
            slides_1 = 0
            for j in range(slides):
                slide_cost = mins_per_slide + (acts_per_slide[j] if j < len(acts_per_slide) else 0)
                if cumulative + slide_cost > remaining_mins:
                    break
                cumulative += slide_cost
                slides_1 = j + 1
            slides_2 = slides - slides_1
            if slides_1 > 0:
                # At least one slide fits today: split across the day boundary
                ref = make_slide_ref(row, slides_1)
                append_first_half(row, slides_1, sum(acts_per_slide[:slides_1]), ref)
                end_day_and_start_next()
                rows_by_deck[i] = make_second_half_row(row, slides_1, slides_2, ref, acts_per_slide[slides_1:])
                continue
            else:
                # Nothing fits today: push whole lesson to next day
                end_day_and_start_next()

        # ---- Schedule lesson ----
        out.append({
            "Sub-lesson": row["Sub-lesson"],
            "Slides": slides,
            "Activity Minutes": activity_mins,
            "kind": "lesson",
        })
        t += timedelta(minutes=total)
        first_lesson_scheduled_today = True

        # ---- Lunch AFTER: lesson end within ±5 min of lunch target ----
        if not lunch_done_today and t >= lunch_earliest:
            if abs(t - lunch_td) <= BREAK_WINDOW:
                append_block("lunch", 45)
                lunch_done_today = True
                skip_post_lunch_break()

        # ---- Break AFTER: lesson end within ±5 min of break target ----
        if next_break_idx < len(targets):
            target = targets[next_break_idx]
            if t >= target - BREAK_WINDOW and t <= target + BREAK_WINDOW:
                if t + timedelta(minutes=10) <= end_td:
                    append_block("break", 10)
                next_break_idx += 1

        i += 1

    # After last lesson, add final Close row
    out.append({
        "Sub-lesson": "Close – Photo, Feedback & Exam",
        "Slides": 0,
        "Activity Minutes": 0,
        "kind": "close",
    })

    return out

# ---------- interactive configuration ----------
def _apply_even_distribution(
    modules: List[str],
    manual_pins: Dict[str, int],
    num_days: int,
    module_weights: Dict[str, float] = None,
) -> Dict[str, int]:
    """Add day boundaries for non-pinned modules.

    When module_weights are supplied, uses time-balanced cuts: finds the
    module boundary whose cumulative minutes is nearest to each day's target,
    with two tie-breaking rules:
      - Complete-module: slightly prefer cutting *after* a module (over target)
        rather than stopping just short, so no module is left dangling.
      - Favour-earlier: processes day boundaries in order so earlier days
        claim their share first; when days are unequal, day 1 is slightly
        longer rather than the last day.

    Falls back to count-based (ceil division) when no weights are given.
    """
    result = dict(manual_pins)

    pin_by_idx = sorted(
        [(modules.index(m), d) for m, d in manual_pins.items() if m in modules],
        key=lambda x: x[1],
    )
    prev_mod_idx, prev_day = 0, 1
    segments = []
    for pin_mod_idx, pin_day in pin_by_idx:
        if pin_mod_idx > prev_mod_idx:
            segments.append((prev_mod_idx, prev_day, pin_mod_idx, pin_day))
        prev_mod_idx, prev_day = pin_mod_idx + 1, pin_day + 1
    if prev_mod_idx < len(modules):
        segments.append((prev_mod_idx, prev_day, len(modules), num_days + 1))

    for mod_start, day_start, mod_end, day_end in segments:
        seg_mods = modules[mod_start:mod_end]
        seg_days = day_end - day_start
        if len(seg_mods) == 0 or seg_days <= 1:
            continue

        if module_weights:
            seg_weights = [module_weights.get(m, 1.0) for m in seg_mods]
            total = sum(seg_weights)
            prev_cut = -1
            for day_offset in range(1, seg_days):
                target_cumul = total * day_offset / seg_days
                best_idx, best_eff = prev_cut, float("inf")
                running = 0.0
                for idx, w in enumerate(seg_weights):
                    running += w
                    if idx <= prev_cut:
                        continue
                    diff = abs(running - target_cumul)
                    # Slightly prefer over-target so the current module is
                    # completed; this also ensures earlier days are favoured
                    eff = diff if running >= target_cumul else diff + 0.01 * total
                    if eff < best_eff:
                        best_eff, best_idx = eff, idx
                if best_idx + 1 < len(seg_mods):
                    result[seg_mods[best_idx + 1]] = day_start + day_offset
                    prev_cut = best_idx
        else:
            per_day = (len(seg_mods) + seg_days - 1) // seg_days
            for day_offset in range(1, seg_days):
                idx = day_offset * per_day
                if idx < len(seg_mods):
                    result[seg_mods[idx]] = day_start + day_offset

    return result


def interactive_config(
    parsed_rows: List[Dict[str, Any]],
    mins_per_slide: float = 2.0,
) -> Dict[str, Any]:
    """Prompt the user for timing and lesson-balancing settings.

    Returns a dict with keys:
        day_start, day_end, lunch_time, day_window, break_targets, lesson_day_map
    """
    print("\n=== SAFe Timing Sheet Configuration ===\n")

    # ── Day Timing ──────────────────────────────────────────────────────────────
    print("── Day Timing ──")
    print("  1. 9:00 – 18:00  (lunch 13:00)  [default]")
    print("  2. 8:30 – 18:00  (lunch 12:30)")
    print("  3. Custom start time")
    timing_choice = input("Choice [1]: ").strip() or "1"

    if timing_choice == "2":
        day_start, day_end, lunch_time = "08:30", "18:00", "12:30"
        day_window = "8:30-18:00"
        cfg_break_targets = None  # use fixed_break_targets
    elif timing_choice == "3":
        while True:
            raw = input("Start time (HH:MM) [09:00]: ").strip() or "09:00"
            if re.match(r'^\d{1,2}:\d{2}$', raw):
                h, m = raw.split(":")
                day_start = f"{int(h):02d}:{m}"
                break
            print("  Please enter time as HH:MM (e.g. 08:00)")
        day_end = "18:00"
        lunch_time = compute_lunch_time(day_start)
        day_window = "custom"
        cfg_break_targets = compute_break_targets(day_start, day_end, lunch_time)
    else:
        day_start, day_end, lunch_time = "09:00", "18:00", "13:00"
        day_window = "9:00-18:00"
        cfg_break_targets = None

    # ── Lesson Balancing ────────────────────────────────────────────────────────
    print("\n── Lesson Balancing ──")
    seen: List[str] = []
    for r in parsed_rows:
        m = r.get("deck_module", "")
        if m and m not in seen:
            seen.append(m)
    modules = seen
    if modules:
        print(f"  Detected modules: {', '.join(modules)}")

    # Step 1: anchors — ask this first so the user can state what they know
    print("\n  Anchor any module to a specific day? (e.g. 7:4 pins module 7 to day 4)")
    print("  Press Enter with no input when done.")
    manual_pins: Dict[str, int] = {}
    while True:
        entry = input("  > ").strip()
        if not entry:
            break
        if ":" in entry:
            parts = entry.split(":", 1)
            if parts[1].strip().isdigit():
                manual_pins[parts[0].strip()] = int(parts[1].strip())
                continue
        print("  Invalid format. Use MODULE:DAY (e.g. 7:4)")

    # Step 2: number of days — default to highest anchored day if available
    default_days = max(manual_pins.values()) if manual_pins else None
    default_hint = str(default_days) if default_days else "auto"
    raw = input(f"\n  Number of days [{default_hint}]: ").strip()
    if raw.isdigit() and int(raw) > 0:
        num_days = int(raw)
    elif not raw and default_days:
        num_days = default_days
    else:
        num_days = None

    # Compute total minutes per module for time-balanced distribution
    module_weights: Dict[str, float] = {}
    for r in parsed_rows:
        m = r.get("deck_module", "")
        if m:
            module_weights[m] = (module_weights.get(m, 0.0)
                                 + r["Slides"] * mins_per_slide
                                 + r["Activity Minutes"])

    # Step 3: distribution for remaining modules
    remaining = [m for m in modules if m not in manual_pins]
    pinned_days = sorted(set(manual_pins.values()))
    if pinned_days and min(pinned_days) > 1:
        remaining_day_range = f"days 1\u2013{min(pinned_days) - 1}"
    elif num_days:
        remaining_day_range = f"days 1\u2013{num_days}"
    else:
        remaining_day_range = "available days"

    dist_choice = "1"
    if remaining:
        print(f"\n  Distribute remaining modules ({', '.join(remaining)}) across {remaining_day_range}:")
        print("    1. Weighted – natural content overflow  [default]")
        print("    2. Balanced – equalise total time per day")
        dist_choice = input("  Choice [1]: ").strip() or "1"

    if dist_choice == "2":
        if num_days is None:
            while True:
                raw = input("  Number of days: ").strip()
                if raw.isdigit() and int(raw) > 0:
                    num_days = int(raw)
                    break
                print("  Please enter a positive integer.")
        lesson_day_map = _apply_even_distribution(
            modules, manual_pins, num_days, module_weights=module_weights
        )
    else:
        lesson_day_map = manual_pins

    # ── Summary ─────────────────────────────────────────────────────────────────
    print("\n── Summary ──")
    print(f"  Start time : {day_start}")
    print(f"  Day window : {day_start} – {day_end}  (lunch {lunch_time})")
    if lesson_day_map or manual_pins:
        pinned_mods = [m for m in modules if m in manual_pins]
        auto_mods   = [m for m in modules if m not in manual_pins and m in lesson_day_map]
        free_mods   = [m for m in modules if m not in lesson_day_map]
        if free_mods:
            free_days = (f"days 1\u2013{min(lesson_day_map.values()) - 1}"
                         if lesson_day_map else f"days 1\u2013{num_days or '?'}")
            dist_label = "weighted" if dist_choice == "1" else "balanced"
            print(f"  Modules {', '.join(free_mods)} \u2192 {free_days}  ({dist_label})")
        for mod in auto_mods:
            print(f"  Module {mod} \u2192 day {lesson_day_map[mod]}  (balanced)")
        for mod in pinned_mods:
            print(f"  Module {mod} \u2192 day {manual_pins[mod]}  (anchored)")
    else:
        print("  Balance    : weighted (natural overflow)")

    while True:
        proceed = input("\nProceed? [y]: ").strip().lower() or "y"
        if proceed in ("y", "yes"):
            break
        elif proceed in ("n", "no"):
            print("Aborted.")
            sys.exit(0)

    return {
        "day_start": day_start,
        "day_end": day_end,
        "lunch_time": lunch_time,
        "day_window": day_window,
        "break_targets": cfg_break_targets,
        "lesson_day_map": lesson_day_map,
    }


# ---------- parallel helper (top-level so it's picklable) ----------
def _parse_deck_task(args):
    path_str, deck_label = args
    rows = parse_deck(Path(path_str), deck_label)
    return path_str, deck_label, rows

# ---------- runner ----------
def main():
    ap = argparse.ArgumentParser(description="SAFe timing → Excel with hard cut-offs and start-of-day rows.")
    ap.add_argument("--input-folder", default="/Users/ecp/Timing-Scripts/decks")
    ap.add_argument("--output", default="TimingSheet.xlsx")
    ap.add_argument("--mins-per-slide", type=float, default=2.0)
    ap.add_argument("--day-window", choices=["8:30-18:00", "9:00-18:00"], default="9:00-18:00",
                    help="Pre-set day window (only used with --no-config)")
    ap.add_argument("--no-open", action="store_true")
    ap.add_argument("--verbose", action="store_true")
    ap.add_argument("--no-config", action="store_true",
                    help="Skip interactive prompts and use CLI flag defaults")
    args = ap.parse_args()

    if args.day_window == "8:30-18:00":
        day_start, day_end, lunch_time = "08:30", "18:00", "12:30"
    else:
        day_start, day_end, lunch_time = "09:00", "18:00", "13:00"

    in_dir = Path(args.input_folder).expanduser().resolve()
    if not in_dir.exists():
        print(f"✗ Folder not found: {in_dir}")
        sys.exit(1)

    pptx_files = sorted([p for p in in_dir.glob("*.pptx") if p.is_file()])
    if not pptx_files:
        print(f"✗ No .pptx files found in {in_dir}")
        sys.exit(1)

    # Parse decks (parallel when multiple files)
    tasks = [(str(f), deck_label_from_filename(f.name)) for f in pptx_files]
    deck_results: Dict[str, Tuple] = {}

    if len(pptx_files) > 1:
        with ProcessPoolExecutor() as executor:
            futures = {executor.submit(_parse_deck_task, t): t for t in tasks}
            for future in as_completed(futures):
                try:
                    path_str, deck_label, rows = future.result()
                    deck_results[deck_label] = (rows, Path(path_str).name)
                    print(f"✓ {Path(path_str).name}  → {len(rows)} rows", flush=True)
                except Exception as e:
                    _, deck_label = futures[future]
                    print(f"✗ {deck_label}  FAILED: {e}", flush=True)
    else:
        for path_str, deck_label in tasks:
            try:
                _, _, rows = _parse_deck_task((path_str, deck_label))
                deck_results[deck_label] = (rows, Path(path_str).name)
                print(f"✓ {Path(path_str).name}  → {len(rows)} rows", flush=True)
            except Exception as e:
                print(f"✗ {Path(path_str).name}  FAILED: {e}", flush=True)

    # Flatten results preserving original file order
    parsed_rows = []
    for f in pptx_files:
        deck_label = deck_label_from_filename(f.name)
        if deck_label not in deck_results:
            continue
        rows, _ = deck_results[deck_label]
        m_mod = re.match(r'^\s*(\d+)', deck_label)
        deck_module = m_mod.group(1) if m_mod else ""
        for name, slides, act, kind, slide_start, acts_per_slide in rows:
            parsed_rows.append(
                {
                    "Deck": deck_label,
                    "Sub-lesson": name,
                    "Slides": int(slides or 0),
                    "Activity Minutes": int(act or 0),
                    "kind": kind,
                    "slide_start": slide_start,
                    "deck_module": deck_module,
                    "acts_per_slide": acts_per_slide,
                }
            )

    if not parsed_rows:
        print("✗ No data parsed.")
        sys.exit(1)

    # Sort by deck then sublesson code
    def subcode(s: str) -> str:
        m = re.match(r"^\s*(\d+\.\d+)", s or "")
        return m.group(1) if m else ""

    parsed_rows.sort(key=lambda r: (r["Deck"], subcode(r["Sub-lesson"]), r["Sub-lesson"]))

    # Interactive configuration (unless suppressed)
    cfg_break_targets = None
    lesson_day_map: Dict[str, int] = {}
    if not args.no_config:
        cfg = interactive_config(parsed_rows, mins_per_slide=args.mins_per_slide)
        day_start    = cfg["day_start"]
        day_end      = cfg["day_end"]
        lunch_time   = cfg["lunch_time"]
        args.day_window = cfg["day_window"]
        cfg_break_targets = cfg["break_targets"]
        lesson_day_map    = cfg["lesson_day_map"]

    # Build sequence with hard cut-offs
    sequence = build_sequence(
        parsed_rows,
        day_start,
        day_end,
        lunch_time,
        args.day_window,
        mins_per_slide=args.mins_per_slide,
        lesson_day_map=lesson_day_map,
        break_targets=cfg_break_targets,
    )

    out_path = Path(args.output).expanduser().resolve()
    tmp_path = out_path.with_suffix(".tmp.xlsx")

    # Build workbook directly with openpyxl (no pandas needed)
    wb = Workbook()
    ws = wb.active
    ws.title = "Timing"
    ws.append(["Sub-lesson", "Slides", "Minutes", "Activity Minutes", "Total Minutes", "Start", "End"])
    for r in sequence:
        ws.append([r["Sub-lesson"], r["Slides"], None, r["Activity Minutes"], None, None, None])

    # Settings sheet
    settings = wb["Settings"] if "Settings" in wb.sheetnames else wb.create_sheet("Settings")
    settings["A1"] = "MinsPerSlide"
    settings["B1"] = float(args.mins_per_slide)
    settings["A2"] = "StartOfDay"
    settings["B2"] = day_start
    settings["A3"] = "EndOfDay"
    settings["B3"] = day_end

    # Apply formulas & formats
    first_data_row = 2
    last_row = ws.max_row
    time_fmt = "h:mm AM/PM"
    grey = PatternFill(fill_type="solid", start_color="EEEEEE", end_color="EEEEEE")

    # Helper: detect row label
    def label(r):
        return (ws[f"A{r}"].value or "").strip()

    for r in range(first_data_row, last_row + 1):
        lbl = label(r)

        # Minutes = Slides * Settings!$B$1 (for all ordinary rows)
        ws[f"C{r}"] = f'=IF(B{r}="","",B{r}*Settings!$B$1)'

        # Total Minutes:
        if lbl == "Close – Photo, Feedback & Exam":
            # Fill remaining time to EndOfDay (never negative)
            ws[f"E{r}"] = f"=MAX(0,(TIMEVALUE(Settings!$B$3)-F{r})*1440)"
        else:
            ws[f"E{r}"] = f'=IF(C{r}="","",C{r}+D{r})'

        # Start / End formulas:
        if r == first_data_row:
            # First row gets StartOfDay
            ws[f"F{r}"] = "=TIMEVALUE(Settings!$B$2)"
        else:
            # If this is a 'Start of Day' marker, reset Start to StartOfDay; otherwise chain from previous End
            if lbl == "Start of Day":
                ws[f"F{r}"] = "=TIMEVALUE(Settings!$B$2)"
            else:
                ws[f"F{r}"] = f"=G{r-1}"

        ws[f"G{r}"] = f"=F{r}+E{r}/1440"

        # Display format for Start/End
        ws[f"F{r}"].number_format = time_fmt
        ws[f"G{r}"].number_format = time_fmt

        # Grey fill for Break/Lunch/Start of Day
        if lbl in ("Break", "Lunch", "Start of Day"):
            for col in ("A", "B", "C", "D", "E", "F", "G"):
                ws[f"{col}{r}"].fill = grey

    # Column widths
    for col, w in {"A": 60, "B": 8, "C": 10, "D": 16, "E": 14, "F": 12, "G": 12}.items():
        ws.column_dimensions[col].width = w

    wb.save(tmp_path)

    # Replace original file
    try:
        if out_path.exists():
            out_path.unlink()
        Path(tmp_path).rename(out_path)
        final_path = out_path
    except Exception:
        ts_path = out_path.with_name(out_path.stem + f".{int(time.time())}.xlsx")
        Path(tmp_path).rename(ts_path)
        final_path = ts_path
        print(f"⚠️ Could not replace {out_path.name} (maybe open). Kept: {final_path}")

    print(f"\n✅ Wrote {last_row-1} rows")
    print(f"📄 File: {final_path}")

    if (not args.no_open) and platform.system() == "Darwin":
        os.system(f'open "{final_path}"')

if __name__ == "__main__":
    main()
